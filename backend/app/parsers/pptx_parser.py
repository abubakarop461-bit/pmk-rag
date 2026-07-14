from pptx import Presentation
from typing import List
from langchain_core.documents import Document
from app.parsers.base_parser import BaseParser
from loguru import logger

class PptxParser(BaseParser):
    def parse(self, file_path: str) -> List[Document]:
        """
        Loads PowerPoint presentation slides and extracts all text box layers.
        """
        logger.info(f"PptxParser starting extraction on: {file_path}")
        try:
            prs = Presentation(file_path)
            documents = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        
                content = f"--- Slide {slide_num + 1} ---\n" + "\n".join(slide_text)
                documents.append(
                    Document(
                        page_content=content.strip(),
                        metadata={
                            "source": file_path,
                            "page": slide_num + 1
                        }
                    )
                )
            logger.info(f"PptxParser completed successfully. Extracted {len(documents)} slides.")
            return documents
        except Exception as e:
            logger.error(f"python-pptx failed to parse PowerPoint slides from {file_path}: {e}")
            raise RuntimeError(f"PowerPoint parsing failed: {e}")
